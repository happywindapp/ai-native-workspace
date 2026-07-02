#!/usr/bin/env python3
"""
Office document extractor — dispatches by file extension.

Extracts xlsx/xls, docx, pdf, pptx to plain text/markdown so Claude can read
content without per-file permission prompts. For scanned/image PDFs or when
visual layout matters, use the `ai-multimodal` skill instead.

Usage:
    py extract.py <file> [<file> ...] [--output PATH] [--max-rows N]

    --output PATH   Write combined result to a file instead of stdout
                    (use for large docs to keep them out of context).
    --max-rows N    Max rows per spreadsheet sheet (default 500).
"""
import sys
from pathlib import Path

# Windows consoles default to cp1252 and choke on Vietnamese text — force UTF-8.
try:
    sys.stdout.reconfigure(encoding="utf-8")
except (AttributeError, ValueError):
    pass


def extract_xlsx(path, max_rows=500):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    out = [f"# FILE: {path}"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out.append(f"\n## SHEET: {sheet_name}")
        rows = 0
        for row in ws.iter_rows(values_only=True):
            vals = list(row)
            while vals and vals[-1] is None:
                vals.pop()
            if not vals:
                continue
            line = " | ".join(str(v) if v is not None else "" for v in vals)
            if line.strip():
                out.append(line)
            rows += 1
            if rows > max_rows:
                out.append(f"... (truncated at {max_rows} rows)")
                break
    return "\n".join(out)


def extract_docx(path):
    from docx import Document
    doc = Document(path)
    out = [f"# FILE: {path}"]
    for p in doc.paragraphs:
        txt = p.text.strip()
        if txt:
            style = p.style.name if p.style else ""
            out.append(f"[{style}] {txt}")
    for ti, table in enumerate(doc.tables):
        out.append(f"\n## TABLE {ti} ({len(table.rows)}x{len(table.columns)})")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " / ") for c in row.cells]
            out.append(" | ".join(cells))
    return "\n".join(out)


def extract_pdf(path):
    from pypdf import PdfReader
    reader = PdfReader(path)
    out = [f"# FILE: {path} ({len(reader.pages)} pages)"]
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        out.append(f"\n## PAGE {i + 1}")
        out.append(text if text else "(no extractable text — likely scanned; use ai-multimodal skill)")
    return "\n".join(out)


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = [f"# FILE: {path} ({len(prs.slides)} slides)"]
    for i, slide in enumerate(prs.slides):
        out.append(f"\n## SLIDE {i + 1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        out.append(txt)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " / ") for c in row.cells]
                    out.append(" | ".join(cells))
    return "\n".join(out)


DISPATCH = {
    ".xlsx": extract_xlsx, ".xlsm": extract_xlsx, ".xls": extract_xlsx,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
}


def main():
    args = sys.argv[1:]
    output, max_rows, files = None, 500, []
    i = 0
    while i < len(args):
        if args[i] == "--output":
            output = args[i + 1]; i += 2
        elif args[i] == "--max-rows":
            max_rows = int(args[i + 1]); i += 2
        else:
            files.append(args[i]); i += 1

    if not files:
        print(__doc__)
        sys.exit(1)

    results = []
    for f in files:
        ext = Path(f).suffix.lower()
        fn = DISPATCH.get(ext)
        if not fn:
            results.append(f"ERROR: unsupported extension '{ext}' for {f} "
                           f"(supported: {', '.join(sorted(DISPATCH))})")
            continue
        try:
            kwargs = {"max_rows": max_rows} if fn is extract_xlsx else {}
            results.append(fn(f, **kwargs))
        except Exception as e:
            results.append(f"ERROR reading {f}: {e}")

    combined = ("\n\n" + "=" * 80 + "\n\n").join(results)
    if output:
        Path(output).write_text(combined, encoding="utf-8")
        print(f"Wrote {len(combined)} chars from {len(files)} file(s) to {output}")
    else:
        print(combined)


if __name__ == "__main__":
    main()